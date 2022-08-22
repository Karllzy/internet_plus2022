#coding: utf-8
"""Images to PPTX file

Designed by L.z.y 2022.06.14 01:27

This script allows the user to convert images to power point(.pptx)

This tool accepts comma image formats (.jpg, .jpeg, .png).

This script requires that `python-ppt` be installed within the Python
environment you are running this script in.

This file can also be imported as a module and contains the following
functions:

    * images2pptx - returns the Presention file
    * main - the main function of the script
"""


from ast import arg, main
from cmath import log
from genericpath import isdir, isfile
from urllib.parse import unquote
import logging
import os
from pptx import Presentation
from pptx.util import Inches
import pathlib
import argparse
import glob
import re


def atoi(text):
    return int(text) if text.isdigit() else text


def natural_keys(text):
    return [atoi(c) for c in re.split(r'(\d+)', text)]


def images2pptx(ppt_img_dir: pathlib.Path, output_file: str) -> Presentation:
    """
    convert images to ppt.
    """
    prs = Presentation()
    prs.slide_height = Inches(7.5)
    prs.slide_width = Inches(13.333)
    blank_layout = prs.slide_layouts[6]
    if isdir(ppt_img_dir):
        files = [glob.glob(os.path.join(ppt_img_dir, suffix))
                for suffix in ["*.png", "*.jpeg", "*.jpg", "*.JPG", "*.PNG", "*.JPEG"]]
        images, comments = [], []
        _ = [images.extend(image_files) for image_files in files]
    elif isfile(ppt_img_dir) and (ppt_img_dir.suffix == '.md') or (ppt_img_dir.suffix == '.markdown'):
        with open(ppt_img_dir, 'r') as f:
            content = f.read()
        splited = re.findall(r'!\[(.*?)\]\((.*?)\)\n\n(.*?)\n\n', content)
        images, comments = [section[1] for section in splited], [section[2] for section in splited]
        print(images, comments)
        ppt_file_path = os.path.split(ppt_img_dir)[0]
        images = [im if isfile(im) else os.path.join(ppt_file_path, im) for im in images]  # 如果是相对路径，则转换为绝对路径
        images = [im if isfile(im) else unquote(im, 'utf-8') for im in images]  # for markdown file
    else:
        logging.error('Input is not a valid dir or file!')
        return
    images = [im for im in images if isfile(im)]  # filter out non-existing images
    images.sort(key=natural_keys)    
    left = top = Inches(0)
    height, width = Inches(7.5), Inches(13.333)
    for idx, img_path in enumerate(images):
        slide = prs.slides.add_slide(blank_layout)
        # add picture to the slide
        slide.shapes.add_picture(
            img_path, left, top, height=height, width=width)
        # add comment to the slide
        if len(comments):
            slide.notes_slide.notes_text_frame.text = comments[idx]
            print(idx, comments[idx])
    prs.save(output_file)
    return prs


def main():
    parser = argparse.ArgumentParser(description="PPT format change")
    parser.add_argument("-i", '--input_dir',type=pathlib.Path, default='./ppt_imgs',
                        help="input ppt dir or markdown file")
    parser.add_argument("-o", '--output_file', default="./ppt_output.pptx",
                        type=pathlib.Path, help='output dir or file')
    args = parser.parse_args()
    images2pptx(args.input_dir, args.output_file)


if __name__ == "__main__":
    main()
